import os
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx.dml.color import RGBColor
import requests
from PIL import Image
import io
from dotenv import load_dotenv
import json
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm

# Load environment variables
load_dotenv()

# Constants
DEFAULT_FONT_SIZE_TITLE = Pt(44)
DEFAULT_FONT_SIZE_SUBTITLE = Pt(24)
DEFAULT_FONT_SIZE_CONTENT = Pt(18)
DEFAULT_FONT_SIZE_SLIDE_TITLE = Pt(36)
DEFAULT_TEXT_COLOR = RGBColor(51, 51, 51)
DEFAULT_IMAGE_HEIGHT = Inches(1.5)  
DEFAULT_IMAGE_WIDTH = Inches(1.5)   
DEFAULT_IMAGE_Y_POSITION = Inches(2.5)  # Adjusted position
DEFAULT_IMAGE_X_POSITION = Inches(3.5)  # Position from right edge


class PPTGenerator:
    def __init__(self, api_key=None):
        """Initialize the PPT Generator with Gemini API"""
        self.api_key = api_key or os.getenv('GEMINI_API_KEY')
        if not self.api_key:
            raise ValueError("Gemini API key is required. Set GEMINI_API_KEY in .env or pass it directly")

        # Configure Gemini
        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel('gemini-2.5-pro')
        self.presentation = Presentation()

    # ---------- Content Generation ----------
    def generate_content_outline(self, topic, num_slides=5):
        """Generate content outline using Gemini"""
        prompt = f"""
        Create a professional PowerPoint outline on "{topic}" with {num_slides} slides.
        
        The presentation should follow this structure:
        1. Title Slide (Introduction)
        2. Overview/Introduction
        3. History/Background
        4. Key Concepts/Principles
        5. Applications/Use Cases
        6. Advantages and Benefits
        7. Disadvantages/Limitations
        8. Current Trends/Recent Developments
        9. Future Prospects
        10. Conclusion/Summary
        
        Return the response as a JSON array with the following structure:
        [
            {{
                "title": "Slide Title",
                "content": "• First detailed bullet point with explanation\n• Second detailed bullet point with explanation\n• Third detailed bullet point with explanation\n• Fourth detailed bullet point with explanation",
                "slide_type": "title|introduction|history|concepts|applications|advantages|disadvantages|trends|future|conclusion",
                "image_needed": true/false,
                "image_description": "Specific image description for this slide"
            }}
        ]

        CRITICAL REQUIREMENTS:
        - The "content" field MUST contain FULL detailed bullet points with explanations, NOT just topic names
        - Each bullet point should be 1-2 sentences explaining the concept
        - Include specific facts, examples, and details for each bullet point
        - Make content educational and informative
        - Use bullet points (•) format
        - Ensure logical flow and progression
        - Include image_needed field for each slide
        - Provide specific image descriptions for slides that need visuals
        
        Example content format:
        "content": "• Machine Learning is a subset of AI that enables computers to learn from data without explicit programming \n• Supervised learning uses labeled training data to make predictions on new, unseen data\n• Unsupervised learning finds hidden patterns in data without predefined labels\n• Deep learning uses neural networks with multiple layers to process complex data"

        The response must be a valid JSON array.
        """

    

        try:
            response = self.model.generate_content(prompt)
            content = response.text.strip()

            if "```json" in content:
                content = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                content = content.split("```")[1].strip()

            return json.loads(content)
        except Exception as e:
            print(f"Error generating content: {e}")
            return self._get_fallback_outline(topic, num_slides)

    def _get_fallback_outline(self, topic, num_slides):
        """Fallback outline if Gemini fails"""
        return [
            {
                "title": f"Introduction to {topic}",
                "content": "• Comprehensive overview of the topic \n• Key objectives and learning goals\n• What to expect from this presentation\n• Importance and relevance in today's world",
                "slide_type": "title",
                "image_needed": True,
                "image_description": f"{topic} concept visualization"
            },
            {
                "title": "Overview and Background",
                "content": "• Definition and core concepts \n• Historical development and evolution\n• Current state and significance\n• Key terminology and definitions",
                "slide_type": "introduction",
                "image_needed": False,
                "image_description": ""
            },
            {
                "title": "History and Development",
                "content": "• Historical timeline and milestones \n• Key inventors and contributors\n• Major breakthroughs and discoveries\n• Evolution over time",
                "slide_type": "history",
                "image_needed": True,
                "image_description": "historical timeline or development chart"
            },
            {
                "title": "Key Concepts and Principles",
                "content": "• Fundamental principles and theories \n• Core concepts and methodologies\n• Important frameworks and models\n• Essential knowledge and understanding",
                "slide_type": "concepts",
                "image_needed": False,
                "image_description": "concept diagram or framework visualization"
            },
            {
                "title": "Applications and Use Cases",
                "content": "• Real-world applications and implementations \n• Industry use cases and examples\n• Practical applications in various fields\n• Case studies and success stories",
                "slide_type": "applications",
                "image_needed": True,
                "image_description": "application examples or use case scenarios"
            },
            {
                "title": "Advantages and Benefits",
                "content": "• Key advantages and positive aspects\n• Benefits and improvements\n• Competitive advantages\n• Value proposition and strengths",
                "slide_type": "advantages",
                "image_needed": False,
                "image_description": ""
            },
            {
                "title": "Disadvantages and Limitations",
                "content": "• Current limitations and challenges\n• Potential drawbacks and concerns\n• Areas for improvement\n• Risk factors and considerations",
                "slide_type": "disadvantages",
                "image_needed": False,
                "image_description": ""
            },
            {
                "title": "Current Trends and Developments",
                "content": "• Latest trends and innovations\n• Recent developments and breakthroughs\n• Current research and advancements\n• Market trends and adoption",
                "slide_type": "trends",
                "image_needed": True,
                "image_description": "trend charts or recent developments"
            },
            {
                "title": "Future Prospects",
                "content": "• Future predictions and forecasts\n• Upcoming developments and opportunities\n• Potential impact and implications\n• Recommendations and next steps",
                "slide_type": "future",
                "image_needed": True,
                "image_description": "future vision or roadmap"
            },
            {
                "title": "Conclusion and Summary",
                "content": "• Summary of key points and insights \n• Final thoughts and recommendations\n• Call to action and next steps\n• Questions and discussion points",
                "slide_type": "conclusion",
                "image_needed": False,
                "image_description": ""
            }
        ]


    def generate_image_description(self, slide_content):
        """Generate image description from slide content"""
        prompt = f"Suggest a relevant image description (8-12 words) for this content:\n{slide_content}\nRequirements:\n- Be specific and descriptive (8-12 words)\n- Focus on the main concept or theme\n- Use professional and technical terms when appropriate\n- Ensure the image would add value to the slide\n- Make it suitable for stock photo search"
        try:
            response = self.model.generate_content(prompt)
            return response.text.strip()
        except:
            return "professional abstract illustration"

    # ---------- Image Handling ----------
    def download_image(self, query, save_path="temp_image.jpg"):
        """Download an image from Pexels API"""
        try:
            pexels_api_key = os.getenv('PEXELS_API_KEY')
            if not pexels_api_key:
                print("⚠️ No Pexels API key found, using placeholder")
                return self._create_placeholder(save_path)

            url = "https://api.pexels.com/v1/search"
            headers = {"Authorization": pexels_api_key}
            params = {"query": query, "per_page": 1, "orientation": "landscape"}

            response = requests.get(url, headers=headers, params=params)
            response.raise_for_status()
            data = response.json()

            if not data.get("photos"):
                return self._create_placeholder(save_path)

            image_url = data["photos"][0]["src"]["original"]
            img_response = requests.get(image_url)
            img_response.raise_for_status()

            with open(save_path, "wb") as f:
                f.write(img_response.content)

            return save_path
        except Exception as e:
            print(f"⚠️ Image download error: {e}")
            return self._create_placeholder(save_path)

    def _create_placeholder(self, save_path):
        """Fallback placeholder image"""
        img = Image.new("RGB", (800, 600), color="#4A90E2")
        img.save(save_path)
        return save_path

    # ---------- Slide Creation ----------
    def create_title_slide(self, title, subtitle=""):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = DEFAULT_FONT_SIZE_TITLE
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        if subtitle:
            slide.placeholders[1].text = subtitle
            slide.placeholders[1].text_frame.paragraphs[0].font.size = DEFAULT_FONT_SIZE_SUBTITLE

    def create_content_slide(self, title, content, image_path=None):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Cm(1.27), Cm(1.27), Cm(22.86), Cm(2.54))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = DEFAULT_FONT_SIZE_SLIDE_TITLE
        title_frame.paragraphs[0].font.bold = True

        # Content (left side)
        content_box = slide.shapes.add_textbox(Cm(1.27), Cm(4.57), Cm(11.43), Cm(12.7))
        content_frame = content_box.text_frame
        content_frame.text = content
        for p in content_frame.paragraphs:
            p.font.size = DEFAULT_FONT_SIZE_CONTENT
            p.font.color.rgb = DEFAULT_TEXT_COLOR

        # Image (right side)
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Cm(21.59), Cm(3.81), height=Cm(13.97))

    def create_comparison_slide(self, title, content):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = DEFAULT_FONT_SIZE_SLIDE_TITLE
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        content_shape = slide.placeholders[1]
        content_shape.text = content
        for p in content_shape.text_frame.paragraphs:
            p.font.size = DEFAULT_FONT_SIZE_CONTENT
            p.font.color.rgb = DEFAULT_TEXT_COLOR

    # ---------- Presentation Generator ----------
    def generate_presentation(self, topic, num_slides=5, output_path="generated_presentation.pptx"):
        if not topic.strip():
            raise ValueError("Topic cannot be empty")
        if num_slides < 1 or num_slides > 20:
            raise ValueError("Slides must be between 1 and 20")

        print(f"📊 Generating presentation on: {topic}")
        outline = self.generate_content_outline(topic, num_slides)

        for i, slide_data in enumerate(outline):
            title = slide_data["title"]
            content = slide_data["content"]
            slide_type = slide_data["slide_type"]
            image_needed = slide_data.get("image_needed", False)
            image_description = slide_data.get("image_description", "")

            print(f"➡️ Creating slide {i+1}: {title}")

            if i == 0 or slide_type == "title":
                self.create_title_slide(title, "Generated by Gemini AI")
            elif slide_type in ["advantages", "disadvantages"]:
                self.create_comparison_slide(title, content)
            else:
                image_path = None
                if image_needed:
                    query = image_description if image_description else self.generate_image_description(content)
                    image_path = self.download_image(query)
                self.create_content_slide(title, content, image_path)

        self.presentation.save(output_path)
        print(f"✅ Presentation saved as: {output_path}")
        return output_path
    
    def create_universal_slide(self, title, content, image_path=None):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Cm(1.27), Cm(0.5), Cm(12), Cm(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Content
        content_box = slide.shapes.add_textbox(Cm(0.80), Cm(4.40), Cm(24.00), Cm(14.7))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # ✅ key line

        # Add text
        content_frame.text = content
        for p in content_frame.paragraphs:
            p.font.size = Pt(18)
            p.space_after = Pt(6)

        # ✅ Shrink if still overflowing
        content_frame.fit_text(max_size=18, min_size=12)

        # Image on right
        if image_path and os.path.exists(image_path):
           pic = slide.shapes.add_picture(image_path, (Cm(16.30), Cm(13.40), Cm(8.00), Cm(.7)))
           pic.lock_aspect_ratio = True

        return slide




# ---------- Run Standalone ----------
if __name__ == "__main__":
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        print("❌ Set GEMINI_API_KEY in .env")
        exit(1)

    generator = PPTGenerator(api_key=api_key)
    generator.generate_presentation("artificial intelligence", num_slides=6, output_path="109.pptx")
